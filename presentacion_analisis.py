# Generador de Presentación PowerPoint para Análisis de Viviendas
# pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def crear_presentacion():
    # Crear presentación
    prs = Presentation()
    
    # Configurar tamaño de diapositiva
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Título
    slide_layout = prs.slide_layouts[0]  # Layout de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Análisis de Precios de Viviendas"
    subtitle.text = "Análisis Exploratorio de Datos y Modelos Predictivos\nDataset: Housing Price Prediction"
    
    # Slide 2: Descripción del Dataset
    slide_layout = prs.slide_layouts[1]  # Layout de contenido
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Descripción del Dataset"
    content.text = """• Dataset: Housing Price Prediction
• Observaciones: 545 viviendas
• Variables: 13 características

Variables Numéricas:
• price: Precio de la vivienda (variable objetivo para regresión)
• area: Área en pies cuadrados
• bedrooms: Número de habitaciones
• bathrooms: Número de baños
• stories: Número de pisos
• parking: Espacios de estacionamiento

Variables Categóricas:
• mainroad: Acceso a carretera principal (yes/no)
• guestroom: Cuarto de huéspedes (yes/no)
• basement: Sótano (yes/no)
• hotwaterheating: Calefacción de agua caliente (yes/no)
• airconditioning: Aire acondicionado (yes/no)
• prefarea: Área preferencial (yes/no)
• furnishingstatus: Estado de amueblado (furnished/semi-furnished/unfurnished)"""
    
    # Slide 3: Análisis de la Variable Precio
    slide_layout = prs.slide_layouts[5]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Análisis de la Variable Precio"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Placeholder para imagen
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = "[INSERTAR IMAGEN: analisis_precios.png]\n\n• Distribución de precios con sesgo hacia la derecha\n• Presencia de valores atípicos en el rango superior\n• Test de Shapiro-Wilk confirma no normalidad\n• Rango de precios: $1,750,000 - $13,300,000"
    
    # Slide 4: Matriz de Correlación
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Matriz de Correlación - Variables Numéricas"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = """[INSERTAR IMAGEN: matriz_correlacion.png]

Correlaciones más significativas con el precio:
• Área: Correlación positiva moderada
• Baños: Correlación positiva
• Habitaciones: Correlación positiva débil
• Estacionamientos: Correlación positiva débil

Las variables numéricas muestran correlaciones esperadas entre sí."""
    
    # Slide 5: Variables Categóricas
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Distribución de Variables Categóricas"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = """[INSERTAR IMAGEN: variables_categoricas.png]

Observaciones clave:
• Mayoría de viviendas tienen acceso a carretera principal
• Distribución equilibrada en aire acondicionado
• Pocas viviendas tienen calefacción de agua caliente
• Estado de amueblado: distribución relativamente equilibrada
• Área preferencial: más viviendas en áreas no preferenciales"""
    
    # Slide 6: Precio vs Variables Categóricas
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Relación Precio vs Variables Categóricas"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = """[INSERTAR IMAGEN: precio_vs_categoricas.png]

Insights importantes:
• Viviendas amuebladas tienden a tener precios más altos
• Aire acondicionado incrementa significativamente el precio
• Áreas preferenciales muestran precios superiores
• Acceso a carretera principal influye positivamente en el precio"""
    
    # Slide 7: Relaciones Variables Numéricas
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Relaciones entre Variables Numéricas y Precio"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = """[INSERTAR IMAGEN: relaciones_numericas.png]

Patrones identificados:
• Relación positiva clara entre área y precio
• Mayor número de baños correlaciona con precios más altos
• Habitaciones: relación positiva pero con variabilidad
• Pisos: influencia moderada en el precio
• Estacionamientos: incremento gradual del precio"""
    
    # Slide 8: Análisis de Componentes Principales
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Análisis de Componentes Principales (PCA)"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = """[INSERTAR IMAGEN: analisis_pca.png]

Resultados del PCA:
• Los primeros 3 componentes explican ~60% de la varianza
• PC1: Principalmente características físicas (área, habitaciones, baños)
• PC2: Características de ubicación y amenidades
• Reducción dimensional efectiva para el modelado"""
    
    # Slide 9: Modelos de Clasificación
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Modelos de Clasificación - Estado de Amueblado"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    img_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    img_frame = img_placeholder.text_frame
    img_frame.text = """[INSERTAR IMAGEN: resultados_modelos.png]

Comparación de Modelos:
• Modelo 1: Todas las variables
• Modelo 2: Selección de características (RFE)

Métricas de rendimiento:
• Accuracy, Precision, Recall, F1-Score
• Matrices de confusión para evaluación detallada
• Importancia de características seleccionadas"""
    
    # Slide 10: Conclusiones
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusiones y Recomendaciones"
    content.text = """Hallazgos Principales:
• El área es el factor más determinante del precio
• Las amenidades (AC, área preferencial) incrementan significativamente el valor
• El estado de amueblado influye en la percepción de precio
• Variables categóricas aportan información valiosa para la clasificación

Recomendaciones:
• Considerar transformaciones para normalizar la distribución de precios
• Incluir variables de interacción en modelos futuros
• Explorar modelos no lineales para capturar relaciones complejas
• Validar resultados con datos externos del mercado inmobiliario

Aplicaciones:
• Tasación automática de propiedades
• Recomendaciones de precios para vendedores
• Análisis de mercado inmobiliario
• Identificación de oportunidades de inversión"""
    
    # Guardar presentación
    prs.save('Analisis_Precios_Viviendas.pptx')
    print("Presentación guardada como 'Analisis_Precios_Viviendas.pptx'")
    print("\nInstrucciones:")
    print("1. Ejecutar 'analisis_python_Housing_price_fixed.py' para generar las imágenes")
    print("2. Insertar manualmente las imágenes en las diapositivas correspondientes:")
    print("   - Slide 3: analisis_precios.png")
    print("   - Slide 4: matriz_correlacion.png") 
    print("   - Slide 5: variables_categoricas.png")
    print("   - Slide 6: precio_vs_categoricas.png")
    print("   - Slide 7: relaciones_numericas.png")
    print("   - Slide 8: analisis_pca.png")
    print("   - Slide 9: resultados_modelos.png")

if __name__ == "__main__":
    crear_presentacion()